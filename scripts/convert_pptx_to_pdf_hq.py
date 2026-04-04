"""
High-quality PPTX to PDF converter.
Extracts PPTX content using python-pptx and renders it with proper fonts.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageDraw, ImageFont
import io

pptx_path = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pptx"
pdf_path = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pdf"

# Font paths
FONT_DIR = "C:/Windows/Fonts/"
FONTS = {
    'arial': FONT_DIR + "arial.ttf",
    'arialbd': FONT_DIR + "arialbd.ttf",
    'calibri': FONT_DIR + "calibri.ttf",
    'calibribd': FONT_DIR + "calibrib.ttf",
    'simsun': FONT_DIR + "simsun.ttc",        # 宋体
    'simhei': FONT_DIR + "simhei.ttf",         # 黑体
    'msyh': FONT_DIR + "msyh.ttc",             # 微软雅黑
}

def load_font(size_pt, bold=False, fallback_simsun=False):
    """Load a font at the given point size, converted to pixels at our DPI."""
    # For Chinese text, use SimSun
    if fallback_simsun:
        font_path = FONTS.get('simsun', FONT_DIR + "simsun.ttc")
        try:
            return ImageFont.truetype(font_path, size_pt, encoding='utf-16-be')
        except:
            return ImageFont.load_default()

    # For Western text
    if bold:
        try:
            return ImageFont.truetype(FONTS['arialbd'], size_pt)
        except:
            try:
                return ImageFont.truetype(FONT_DIR + "arialbd.ttf", size_pt)
            except:
                return ImageFont.load_default()
    else:
        try:
            return ImageFont.truetype(FONTS['arial'], size_pt)
        except:
            try:
                return ImageFont.truetype(FONT_DIR + "arial.ttf", size_pt)
            except:
                return ImageFont.load_default()

def has_chinese(text):
    """Check if text contains CJK characters."""
    for c in text:
        if '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f':
            return True
    return False

def rgb_to_tuple(rgb):
    if rgb is None:
        return (255, 255, 255)
    try:
        return (rgb[0], rgb[1], rgb[2])
    except:
        return (255, 255, 255)

def emu_to_px(emu, dpi=200):
    return int(emu / 914400 * dpi)

def draw_gradient_bg(draw, W, H):
    """Draw a deep navy gradient background."""
    # Dark navy: #0A1428
    # Slightly lighter at bottom
    for y in range(H):
        ratio = y / H
        r = int(10 + ratio * 5)
        g = int(20 + ratio * 8)
        b = int(40 + ratio * 15)
        draw.line([(0, y), (W, y)], fill=(r, g, b))

def get_run_color(run, default=(255, 255, 255)):
    """Get the color from a text run."""
    try:
        if run.font and run.font.color and run.font.color.rgb:
            return rgb_to_tuple(run.font.color.rgb)
    except:
        pass
    return default

def get_run_font_size(run, default=14):
    """Get font size from a run."""
    try:
        if run.font and run.font.size:
            return int(run.font.size.pt)
    except:
        pass
    return default

def get_run_bold(run, default=False):
    """Check if run is bold."""
    try:
        if run.font and run.font.bold:
            return True
    except:
        pass
    return default

def get_para_align(para, default=PP_ALIGN.LEFT):
    """Get paragraph alignment."""
    try:
        return para.alignment
    except:
        return default

def is_chinese_char(c):
    return '\u4e00' <= c <= '\u9fff' or '\u3000' <= c <= '\u303f' or '\uff00' <= c <= '\uffef'

def contains_chinese(text):
    return any(is_chinese_char(c) for c in text)

def wrap_text_pil(text, font, max_width):
    """Wrap text to fit max_width pixels."""
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        try:
            bbox = font.getbbox(test)
            w = bbox[2] - bbox[0]
        except:
            bbox = (0, 0, len(test) * font.size * 0.6, font.size)
            w = bbox[2] - bbox[0]
        if w <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    if not lines and text:
        # Single word that's too long, just force it
        lines = [text[:int(max_width / (font.size * 0.6))] + "..."]
    return lines

def draw_rounded_rect(draw, xy, radius, fill=None, outline=None, width=1):
    """Draw a rounded rectangle."""
    x1, y1, x2, y2 = xy
    if fill:
        # Simple rounded rect by drawing rect with corners cut
        draw.rectangle([x1+radius, y1, x2-radius, y2], fill=fill)
        draw.rectangle([x1, y1+radius, x2, y2-radius], fill=fill)
        draw.pieslice([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=fill)
        draw.pieslice([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=fill)
        draw.pieslice([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=fill)
        draw.pieslice([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=fill)
    if outline:
        # Draw outline arcs
        draw.arc([x1, y1, x1+2*radius, y1+2*radius], 180, 270, fill=outline, width=width)
        draw.arc([x2-2*radius, y1, x2, y1+2*radius], 270, 360, fill=outline, width=width)
        draw.arc([x1, y2-2*radius, x1+2*radius, y2], 90, 180, fill=outline, width=width)
        draw.arc([x2-2*radius, y2-2*radius, x2, y2], 0, 90, fill=outline, width=width)
        draw.line([(x1+radius, y1), (x2-radius, y1)], fill=outline, width=width)
        draw.line([(x2, y1+radius), (x2, y2-radius)], fill=outline, width=width)
        draw.line([(x1+radius, y2), (x2-radius, y2)], fill=outline, width=width)
        draw.line([(x1, y1+radius), (x1, y2-radius)], fill=outline, width=width)

def draw_slide_from_pptx(slide, slide_num, total, DPI=200, prs=None):
    """Render a PPTX slide as a PIL Image with high quality."""
    if prs is None:
        prs = slide.part._part

    # Get slide dimensions in EMU
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height
    W = int(slide_w_emu / 914400 * DPI)
    H = int(slide_h_emu / 914400 * DPI)

    img = Image.new('RGBA', (W, H), (10, 20, 40, 255))
    draw = ImageDraw.Draw(img)

    # Draw gradient background
    for y in range(H):
        ratio = y / H
        r = int(10 + ratio * 5)
        g = int(20 + ratio * 10)
        b = int(40 + ratio * 20)
        img.putpixel((0, y), (r, g, b, 255))

    # Color constants
    ACCENT_BLUE = (0, 180, 216)    # #00B4D8
    ACCENT_GOLD = (255, 193, 7)    # #FFC107
    ACCENT_DARK = (45, 74, 97)     # #2D4A61
    NAVY_DARK = (13, 27, 42)       # #0D1B2A
    WHITE = (255, 255, 255)
    GRAY = (180, 190, 210)
    CARD_BG = (20, 40, 80, 200)
    TEXT_MUTED = (140, 155, 180)

    # Collect all shape data
    shapes = []
    for shape in slide.shapes:
        try:
            left = emu_to_px(shape.left, DPI)
            top = emu_to_px(shape.top, DPI)
            width = emu_to_px(shape.width, DPI)
            height = emu_to_px(shape.height, DPI)
        except:
            continue

        shapes.append({
            'left': left, 'top': top,
            'width': width, 'height': height,
            'shape': shape
        })

    # Sort: draw rectangles first, then text
    # Draw background fills (solid color rectangles)
    for sd in shapes:
        shape = sd['shape']
        left, top, width, height = sd['left'], sd['top'], sd['width'], sd['height']

        if not shape.has_text_frame:
            # It's a pure shape - draw it
            try:
                fill = shape.fill
                if fill and hasattr(fill, 'type'):
                    fill_type = fill.type
                    if hasattr(fill_type, 'name'):
                        fname = fill_type.name
                    else:
                        fname = str(fill_type)
                else:
                    fname = 'none'
            except:
                fname = 'none'

            # Draw solid fill rectangles
            try:
                if fill and hasattr(fill, 'fore_color') and fill.fore_color:
                    fc = fill.fore_color
                    if hasattr(fc, 'rgb') and fc.rgb:
                        rgb = fc.rgb
                        color = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
                        # Draw with alpha blending
                        bg_pixel = (10, 20, 40)
                        blended = tuple(int(bg_pixel[i] * (1 - 0.8) + color[i] * 0.8) for i in range(3))
                        draw.rectangle([left, top, left+width, top+height], fill=blended)
                    elif hasattr(fc, 'theme') and fc.theme is not None:
                        # Theme color - use a dark fill
                        draw.rectangle([left, top, left+width, top+height], fill=(20, 35, 65))
            except:
                pass

    # Draw shapes with line/stroke
    for sd in shapes:
        shape = sd['shape']
        left, top, width, height = sd['left'], sd['top'], sd['width'], sd['height']

        try:
            line = shape.line
            if line and hasattr(line, 'fill') and line.fill:
                lc = line.fill
                if hasattr(lc, 'fore_color') and lc.fore_color and hasattr(lc.fore_color, 'rgb'):
                    lrgb = lc.fore_color.rgb
                    line_color = (int(lrgb[0]), int(lrgb[1]), int(lrgb[2]))
                    line_width = emu_to_px(shape.line.width, DPI) if shape.line.width else 1
                    draw.rectangle([left, top, left+width, top+height], outline=line_color, width=max(line_width, 1))
        except:
            pass

    # Draw text content
    for sd in shapes:
        shape = sd['shape']
        if not shape.has_text_frame:
            continue

        left, top, width, height = sd['left'], sd['top'], sd['width'], sd['height']
        tf = shape.text_frame

        # Get shape fill for background (if any)
        try:
            shape_fill = None
            if hasattr(shape, 'fill') and shape.fill and hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                if hasattr(shape.fill.fore_color, 'rgb') and shape.fill.fore_color.rgb:
                    frgb = shape.fill.fore_color.rgb
                    shape_fill = (int(frgb[0]), int(frgb[1]), int(frgb[2]), 220)
        except:
            shape_fill = None

        # Draw text background if shape has fill
        if shape_fill and shape_fill[3] > 0:
            draw.rounded_rectangle([left, top, left+width, top+height], radius=8,
                                   fill=shape_fill[:3])

        # Draw text paragraphs
        y_offset = top
        margin_left = left + 8
        margin_right = left + width - 8

        for para in tf.paragraphs:
            text = para.text.strip()
            if not text:
                y_offset += int(14 * DPI / 72 * 1.3)
                continue

            # Get paragraph alignment
            try:
                align = para.alignment
                align_left = (align == PP_ALIGN.LEFT or align is None)
                align_center = (align == PP_ALIGN.CENTER)
                align_right = (align == PP_ALIGN.RIGHT)
            except:
                align_left = True
                align_center = False
                align_right = False

            # Process runs within paragraph
            if para.runs:
                # Render each run separately
                x_cursor = margin_left
                max_y = y_offset
                for run in para.runs:
                    run_text = run.text
                    if not run_text:
                        continue

                    font_size = get_run_font_size(run, 14)
                    bold = get_run_bold(run, False)
                    color = get_run_color(run, WHITE)
                    is_ch = contains_chinese(run_text)

                    # Choose font
                    if is_ch:
                        try:
                            font = ImageFont.truetype(FONTS['msyh'], int(font_size * DPI / 72), encoding='utf-16-be')
                        except:
                            try:
                                font = ImageFont.truetype(FONTS['simhei'], int(font_size * DPI / 72), encoding='utf-16-be')
                            except:
                                font = load_font(font_size, False)
                    else:
                        font = load_font(font_size, bold)

                    # Wrap this run's text
                    avail_width = margin_right - x_cursor
                    if avail_width <= 0:
                        break
                    wrapped = wrap_text_pil(run_text, font, avail_width)

                    rx = x_cursor
                    for wline in wrapped:
                        if y_offset + font.size > top + height - 4:
                            break
                        ry = y_offset
                        draw.text((rx, ry), wline, font=font, fill=color)
                        try:
                            bbox = font.getbbox(wline)
                            lh = max(bbox[3] - bbox[1], int(font_size * 1.3))
                        except:
                            lh = int(font_size * 1.3)
                        y_offset += lh
                        max_y = max(max_y, y_offset)

                    x_cursor = margin_left  # Reset for next run on new line
                    if len(wrapped) > 1:
                        y_offset = max_y  # Stay at the last line position
            else:
                # No runs - render whole paragraph as one text block
                font_size = 14
                try:
                    if tf.paragraphs[0] and tf.paragraphs[0].runs and tf.paragraphs[0].runs[0]:
                        font_size = get_run_font_size(tf.paragraphs[0].runs[0], 14)
                except:
                    pass

                is_ch = contains_chinese(text)
                font = (ImageFont.truetype(FONTS['msyh'], int(font_size * DPI / 72), encoding='utf-16-be')
                        if is_ch else load_font(font_size, False))

                wrapped = wrap_text_for_para(text, font, width - 16, font_size)
                for wline in wrapped:
                    if y_offset + font_size > top + height - 4:
                        break
                    rx = margin_left
                    draw.text((rx, y_offset), wline, font=font, fill=WHITE)
                    y_offset += int(font_size * 1.3)

        # Also draw the raw text overlay to catch any missed content
        raw_text = tf.text.strip()
        if raw_text:
            # Check if there's Chinese
            if contains_chinese(raw_text):
                try:
                    fnt = ImageFont.truetype(FONTS['msyh'], int(14 * DPI / 72), encoding='utf-16-be')
                except:
                    try:
                        fnt = ImageFont.truetype(FONTS['simhei'], int(14 * DPI / 72), encoding='utf-16-be')
                    except:
                        fnt = load_font(14, False)
            else:
                fnt = load_font(14, False)
            draw.text((margin_left, top + 4), raw_text[:100], font=fnt, fill=WHITE)

    # Footer: slide number
    font_small = load_font(10, False)
    draw.text((W - 80, H - 30), f"{slide_num} / {total}",
               font=font_small, fill=TEXT_MUTED)

    return img

def wrap_text_for_para(text, font, max_width, font_size):
    """Wrap paragraph text."""
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        try:
            bbox = font.getbbox(test)
            w = bbox[2] - bbox[0]
        except:
            w = len(test) * font_size * 0.6
        if w <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    if not lines:
        lines = [text[:int(max_width / (font_size * 0.6))]]
    return lines

# Main
print(f"PPTX: {pptx_path}")
prs = Presentation(pptx_path)
print(f"Slides: {len(prs.slides)}")

W_in = prs.slide_width.inches
H_in = prs.slide_height.inches
DPI = 200  # Higher DPI for better quality

W = int(W_in * DPI)
H = int(H_in * DPI)
print(f"Slide size: {W_in:.2f} x {H_in:.2f} inches = {W} x {H} px @ {DPI} DPI")

print("Rendering slides...")
images = []
total = len(prs.slides)
for i, slide in enumerate(prs.slides, 1):
    print(f"  Slide {i}/{total}...", end="", flush=True)
    img = draw_slide_from_pptx(slide, i, total, DPI, prs)
    images.append(img.convert('RGB'))
    print(" OK")

print(f"\nSaving PDF: {pdf_path}")
images[0].save(
    pdf_path,
    save_all=True,
    append_images=images[1:],
    resolution=DPI
)

size_mb = os.path.getsize(pdf_path) / 1024 / 1024
print(f"Done! PDF: {size_mb:.2f} MB")
