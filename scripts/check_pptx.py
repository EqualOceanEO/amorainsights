from pptx import Presentation
import io

prs = Presentation(r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pptx")
lines = []
lines.append(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    lines.append(f"=== Slide {i+1} ===")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            try:
                txt = shape.text.strip()[:100]
                txt = txt.encode('ascii', 'replace').decode('ascii')
                lines.append(f"  {shape.name}: {txt}")
            except:
                pass

output = "\n".join(lines)
with open(r"C:\Users\51229\WorkBuddy\Claw\scripts\pptx_content.txt", "w", encoding="utf-8") as f:
    f.write(output)
print("Done - written to pptx_content.txt")
