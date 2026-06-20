"""
Convert PPTX to PDF using python-pptx + Pillow.
Extracts slide content and renders as images in a PDF.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import io

pptx_path = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pptx"
pdf_path = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pdf"

print(f"PPTX: {pptx_path}")
print(f"PDF: {pdf_path}")

prs = Presentation(pptx_path)
slide_w_in = prs.slide_width.inches
slide_h_in = prs.slide_height.inches
print(f"Slide size: {slide_w_in:.2f} x {slide_h_in:.2f} inches")

# DPI and pixel dimensions
DPI = 150
W = int(slide_w_in * DPI)
H = int(slide_h_in * DPI)
print(f"Image size: {W} x {H} px @ {DPI} DPI")

# Color palette
BG_DARK = (10, 20, 40)
BG_MID = (15, 30, 60)
BG_CARD = (25, 40, 80)
TEXT_WHITE = (255, 255, 255)
TEXT_GRAY = (180, 190, 210)
ACCENT = (0, 180, 216)
GOLD = (255, 193, 7)
GREEN = (76, 175, 80)
RED = (244, 67, 54)
ORANGE = (255, 152, 0)
PURPLE = (156, 39, 176)

def make_font(size_pt, bold=False):
    try:
        if bold:
            return ImageFont.truetype("C:\\Windows\\Fonts\\arialbd.ttf", int(size_pt * DPI / 72))
        return ImageFont.truetype("C:\\Windows\\Fonts\\arial.ttf", int(size_pt * DPI / 72))
    except:
        return ImageFont.load_default()

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def draw_bg(draw, W, H):
    """Draw gradient background."""
    for y in range(H):
        ratio = y / H
        r = int(BG_DARK[0] + (BG_MID[0] - BG_DARK[0]) * ratio * 0.5)
        g = int(BG_DARK[1] + (BG_MID[1] - BG_DARK[1]) * ratio * 0.5)
        b = int(BG_DARK[2] + (BG_MID[2] - BG_DARK[2]) * ratio * 0.5)
        draw.line([(0, y), (W, y)], fill=(r, g, b))

def draw_shape(shape, draw):
    """Draw a single shape."""
    try:
        left = int(shape.left / 914400 * DPI)
        top = int(shape.top / 914400 * DPI)
        width = int(shape.width / 914400 * DPI)
        height = int(shape.height / 914400 * DPI)
    except:
        return

    # Get fill color
    try:
        fill = shape.fill
        if fill and hasattr(fill, 'fore_color') and fill.fore_color:
            fc = fill.fore_color
            if hasattr(fc, 'rgb') and fc.rgb:
                rgb = fc.rgb
                color = (rgb[0], rgb[1], rgb[2])
                draw.rectangle([left, top, left+width, top+height], fill=color)
    except:
        pass

def get_text_info(para, run=None):
    """Extract text formatting from paragraph/run."""
    font_size = 14
    bold = False
    italic = False
    color = TEXT_WHITE
    font_name = None

    if para and hasattr(para, 'font'):
        f = para.font
        if f.size:
            font_size = f.size.pt
        if f.bold:
            bold = True
        if f.italic:
            italic = True
        if f.color and hasattr(f.color, 'rgb') and f.color.rgb:
            rgb = f.color.rgb
            color = (rgb[0], rgb[1], rgb[2])
        if f.name:
            font_name = f.name

    return font_size, bold, italic, color, font_name

def draw_textbox(shape, draw):
    """Draw a textbox shape."""
    try:
        left = int(shape.left / 914400 * DPI)
        top = int(shape.top / 914400 * DPI)
        width = int(shape.width / 914400 * DPI)
        height = int(shape.height / 914400 * DPI)
    except:
        return

    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    if not paragraphs:
        return

    # Get default paragraph formatting
    default_size, default_bold, default_italic, default_color, default_font = get_text_info(paragraphs[0])

    # Calculate line height
    line_h = max(default_size * 1.2, 20)
    y = top

    for para in paragraphs:
        if not para.text.strip():
            y += line_h * 0.5
            continue

        font_size, bold, italic, color, font_name = get_text_info(para)

        # Choose font
        font_path = None
        try:
            if font_name and 'bold' in font_name.lower():
                font_path = "C:\\Windows\\Fonts\\arialbd.ttf"
            elif bold:
                font_path = "C:\\Windows\\Fonts\\arialbd.ttf"
            else:
                font_path = "C:\\Windows\\Fonts\\arial.ttf"
            font = ImageFont.truetype(font_path, int(font_size * DPI / 72))
        except:
            font = make_font(font_size, bold)

        # Wrap text
        words = para.text.split()
        lines = []
        current = ""
        for word in words:
            test = (current + " " + word).strip()
            try:
                bbox = draw.textbbox((0, 0), test, font=font)
                if bbox[2] - bbox[0] <= width - 10:
                    current = test
                else:
                    if current:
                        lines.append(current)
                    current = word
            except:
                if len(current) < 30:
                    current = test
                else:
                    lines.append(current)
                    current = word
        if current:
            lines.append(current)

        for line in lines:
            if y < top + height - 5 and y < H - 10:
                draw.text((left, y), line, font=font, fill=color)
                try:
                    bbox = draw.textbbox((0, 0), line, font=font)
                    line_height = bbox[3] - bbox[1] + 4
                except:
                    line_height = font_size + 4
                y += max(line_height, font_size * 1.2)

def draw_slide_image(slide, slide_num, total):
    """Render a slide as a PIL Image."""
    img = Image.new('RGB', (W, H), color=BG_DARK)
    draw = ImageDraw.Draw(img)

    draw_bg(draw, W, H)

    # Draw shapes (rectangles first, then text)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            draw_shape(shape, draw)

    # Draw textboxes
    for shape in slide.shapes:
        if shape.has_text_frame:
            draw_textbox(shape, draw)

    # Slide number footer
    font_small = make_font(10)
    draw.text((W - 80, H - 30), f"{slide_num} / {total}",
               font=font_small, fill=(120, 130, 160))

    return img

print("Generating slide images...")
images = []
total = len(prs.slides)
for i, slide in enumerate(prs.slides, 1):
    print(f"  Slide {i}/{total}...", end="", flush=True)
    img = draw_slide_image(slide, i, total)
    images.append(img)
    print(" OK")

print(f"\nSaving PDF: {pdf_path}")
images[0].save(
    pdf_path,
    save_all=True,
    append_images=images[1:],
    resolution=DPI
)

size_mb = os.path.getsize(pdf_path) / 1024 / 1024
print(f"Done! PDF size: {size_mb:.2f} MB")
