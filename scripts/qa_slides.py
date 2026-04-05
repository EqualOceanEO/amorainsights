"""QA check for new Mapping slides"""
from pptx import Presentation
import os

prs = Presentation(r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v3.pptx")
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    texts = []
    shapes_with_text = 0
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            shapes_with_text += 1
            t = shape.text.strip()[:80]
            try:
                t = t.encode('ascii', 'replace').decode('ascii')
            except:
                t = t[:80]
            texts.append(t)

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    print(f"\n=== Slide {slide_num} ({slide_w.inches:.2f}\" x {slide_h.inches:.2f}\") ===")
    print(f"  Shapes with text: {shapes_with_text}")
    for t in texts[:6]:
        print(f"  {t}")

    problems = []
    for shape in slide.shapes:
        if hasattr(shape, 'left') and hasattr(shape, 'top'):
            l = float(shape.left.inches) if shape.left else 0
            t2 = float(shape.top.inches) if shape.top else 0
            w = float(shape.width.inches) if shape.width else 0
            h = float(shape.height.inches) if shape.height else 0
            sw = float(prs.slide_width.inches)
            sh = float(prs.slide_height.inches)
            if l + w > sw + 0.3:
                problems.append(f"Right overflow: x={l:.2f}+w={w:.2f}>{sw:.2f}")
            if t2 + h > sh + 0.3:
                problems.append(f"Bottom overflow: y={t2:.2f}+h={h:.2f}>{sh:.2f}")

    if problems:
        print(f"  PROBLEMS: {'; '.join(problems)}")
    else:
        print(f"  OK")
