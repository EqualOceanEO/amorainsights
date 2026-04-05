"""
Add 5 rich Mapping slides to HRI-2026-Report-Presentation-v2.pptx after slide 8
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

INPUT_PATH = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v2.pptx"
OUTPUT_PATH = r"C:\Users\51229\WorkBuddy\Claw\HRI-2026-Report-Presentation-v3.pptx"

BG_DARK      = RGBColor(0x0A, 0x0A, 0x14)
BG_CARD      = RGBColor(0x1A, 0x1A, 0x32)
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)
ACCENT_PINK  = RGBColor(0xFF, 0x00, 0x6E)
ACCENT_GREEN = RGBColor(0x00, 0xFF, 0x88)
ACCENT_AMBER = RGBColor(0xFF, 0xBE, 0x0B)
ACCENT_PURPLE= RGBColor(0x7B, 0x2C, 0xBF)
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_GRAY    = RGBColor(0x88, 0x88, 0x99)
TEXT_LGRAY   = RGBColor(0xCC, 0xCC, 0xDD)
CHINA_RED    = RGBColor(0xFF, 0x6B, 0x6B)
USA_BLUE     = RGBColor(0x4D, 0xAB, 0xF7)
EU_PURPLE    = RGBColor(0x97, 0x75, 0xFA)

FOOTER_TEXT = "AMORA Insights  |  Humanoid Robot Index 2026  |  Data Cut-off: Q1 2026"

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_pt=0.5):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=13, color=None, bold=False,
             align=PP_ALIGN.LEFT, italic=False):
    if color is None:
        color = TEXT_WHITE
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return tb

def set_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_footer(slide, pn, total=27):
    add_text(slide, FOOTER_TEXT, 0.3, 5.32, 8.5, 0.22, size=8.5, color=TEXT_GRAY)
    add_text(slide, f"{pn} / {total}", 9.0, 5.32, 0.9, 0.22, size=8.5, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)

def add_header(slide, label, title, sub, pn, total=27):
    add_text(slide, label, 0.3, 0.2, 9.4, 0.32, size=10.5, color=ACCENT_CYAN, bold=True)
    add_text(slide, title, 0.3, 0.5, 9.4, 0.55, size=25, color=TEXT_WHITE, bold=True)
    add_text(slide, sub, 0.3, 1.03, 9.4, 0.28, size=12, color=TEXT_GRAY)
    div = add_rect(slide, 0.3, 1.27, 9.4, 0.015, fill_rgb=RGBColor(0x30, 0x30, 0x50))
    add_footer(slide, pn, total)

prs = Presentation(INPUT_PATH)
# Find a usable blank layout
blank = None
for i in range(len(prs.slide_layouts)-1, -1, -1):
    try:
        test = prs.slide_layouts[i]
        blank = test
        print(f"Using slide_layouts[{i}]: {test.name}")
        break
    except:
        pass
if blank is None:
    blank = prs.slide_layouts[0]
total_orig = len(prs.slides)
print(f"Original slides: {total_orig}")

# ===================== SLIDE M3: China Supply Chain =====================
def build_m3(prs):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    add_header(slide, "M - Mapping  China Supply Chain Ecosystem",
               "China's Supply Chain Atlas",
               "30+ Domestic Suppliers  |  Pearl River & Yangtze Delta  |  Full Actuator Ecosystem",
               9, 27)

    cats = [
        ("Vision System", ACCENT_CYAN,
         [("Hesai (合眸)", "AT128 LiDAR", "Listed"), ("RoboSense (速腾)", "M1/M2 LiDAR", "Listed"), ("Orbbec (奥比中光)", "3D Depth Camera", "Listed"), ("DJI Livox", "Horizon/Avia LiDAR", "Private")],
         0, 0),
        ("Compute & AI", ACCENT_PURPLE,
         [("Horizon Robotics (地平线)", "Journey 5/6 SoC", "Listed"), ("Black Sesame (黑芝麻)", "A1000 Chip", "Listed"), ("CATL (宁德时代)", "Robot Battery Pack", "Listed"), ("EVE Energy (亿纬)", "Cylindrical Cells", "Listed")],
         1, 0),
        ("Actuators", ACCENT_GREEN,
         [("Leaderdrive (绿的谐波)", "Harmonic Reducer", "Listed"), ("Estun (埃斯顿)", "Servo Motors", "Listed"), ("Inovance (汇川)", "Servo Systems", "Listed"), ("Welling (威灵)", "Precision Motors", "Private")],
         0, 1),
        ("Dexterous Hand", ACCENT_AMBER,
         [("Inspire Robots (因时)", "RH56DFX Hand", "Private"), ("Paxini (帕西米)", "Tactile Sensors", "Private"), ("DH-Robotics (大寰)", "Electric Grippers", "Private")],
         1, 1),
        ("Force Sensing", ACCENT_PINK,
         [("Unitree (宇树)", "Proprietary F/T Sensor", "IPO 2026"), ("Tacter (坤维)", "6-Axis F/T Sensor", "Private")],
         0, 2),
    ]

    # Card positions: [col_x, cy]
    # col 0: cx=0.2; col 1: cx=5.0
    # row 0: cy=1.45; row 1: cy=3.22; row 2: cy=4.99
    # Force Sensing (col=0, row=2): 2 items fits within 0.63 available
    CARD_COL_X = [0.2, 5.0]
    CARD_CY = [1.45, 3.22, 4.99]

    for cat_name, cat_color, items, col, row in cats:
        cx = CARD_COL_X[col]
        cy = CARD_CY[row]

        # Category header
        hb = add_rect(slide, cx, cy, 4.65, 0.28,
                     fill_rgb=RGBColor(int(cat_color[0]*0.12), int(cat_color[1]*0.12), int(cat_color[2]*0.12)),
                     line_rgb=cat_color, line_pt=0.75)
        add_text(slide, cat_name, cx+0.1, cy+0.04, 4.4, 0.22, size=10.5, color=cat_color, bold=True)

        for si, (s_name, s_prod, s_tag) in enumerate(items):
            sy = cy + 0.32 + si * 0.3
            sb = add_rect(slide, cx, sy, 4.65, 0.27, fill_rgb=BG_CARD,
                         line_rgb=RGBColor(0x30, 0x30, 0x48), line_pt=0.3)
            add_text(slide, s_name, cx+0.1, sy+0.04, 2.5, 0.21, size=9, color=TEXT_WHITE, bold=True)
            add_text(slide, s_prod, cx+2.6, sy+0.04, 1.4, 0.21, size=8.5, color=TEXT_GRAY)
            tag_c = CHINA_RED if s_tag == "Listed" else (ACCENT_AMBER if "IPO" in s_tag else TEXT_GRAY)
            add_text(slide, s_tag, cx+4.1, sy+0.05, 0.5, 0.2, size=7.5, color=tag_c,
                    bold=(s_tag=="Listed"), align=PP_ALIGN.CENTER)

    # Right: OEM list + strengths
    rx = 5.0
    oem_y = 1.45
    add_rect(slide, rx, oem_y, 4.65, 0.28,
             fill_rgb=RGBColor(0x0A, 0x28, 0x14), line_rgb=ACCENT_GREEN, line_pt=0.75)
    add_text(slide, "Leading OEMs", rx+0.12, oem_y+0.04, 4.4, 0.22, size=10.5, color=ACCENT_GREEN, bold=True)

    oems = [
        ("Unitree (宇树科技)", "5,500+ units 2024  |  Global #1", ACCENT_GREEN),
        ("AgiBot (智元机器人)", "Series B+  |  Domestic leader", ACCENT_CYAN),
        ("Fourier Intelligence (傅里叶)", "GR-1/GR-2  |  Medical+Industry", TEXT_LGRAY),
        ("Leju Robotics (乐聚)", "Robovie series  |  Tier-1", TEXT_LGRAY),
        ("UBTECH (优必选)", "Walker X/S  |  HKEX Listed", CHINA_RED),
        ("AiDimension (星动纪元)", "X1 series  |  Series A 2024", TEXT_GRAY),
    ]
    for oi, (name, desc, oc) in enumerate(oems):
        oy = oem_y + 0.32 + oi * 0.29
        ob = add_rect(slide, rx, oy, 4.65, 0.27, fill_rgb=BG_CARD,
                     line_rgb=RGBColor(0x30, 0x30, 0x48), line_pt=0.3)
        add_text(slide, name, rx+0.12, oy+0.04, 2.2, 0.21, size=9, color=oc, bold=True)
        add_text(slide, desc, rx+2.32, oy+0.04, 2.22, 0.21, size=8, color=TEXT_GRAY)

    str_y = oem_y + 0.32 + len(oems)*0.29 + 0.12
    sb2 = add_rect(slide, rx, str_y, 4.65, 1.22,
                  fill_rgb=RGBColor(0x10, 0x10, 0x28), line_rgb=ACCENT_CYAN, line_pt=0.75)
    add_text(slide, "Core Advantages", rx+0.15, str_y+0.06, 4.3, 0.22, size=10.5, color=ACCENT_CYAN, bold=True)
    strengths = [
        ("Manufacturing Density:", "Full ecosystem within 100km; Pearl River + Yangtze Delta cluster"),
        ("Cost Leadership:", "35% of US cost at 100K scale - 65% structural advantage"),
        ("Actuator Mastery:", "30+ domestic suppliers, most complete globally"),
        ("Speed Advantage:", "Unitree 5,500+ deployed vs Tesla Optimus <50"),
    ]
    for si2, (lb, desc2) in enumerate(strengths):
        sy2 = str_y + 0.32 + si2 * 0.21
        add_text(slide, lb, rx+0.18, sy2, 1.52, 0.2, size=9, color=ACCENT_AMBER, bold=True)
        add_text(slide, desc2, rx+1.7, sy2, 2.82, 0.2, size=9, color=TEXT_LGRAY)

    return slide


# ===================== SLIDE M4: US Supply Chain =====================
def build_m4(prs):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    add_header(slide, "M - Mapping  US Supply Chain Ecosystem",
               "US Supply Chain Atlas",
               "NVIDIA CUDA Platform Moat  |  70% Critical Component Import Dependency  |  $2B+ VC Funding 2024",
               10, 27)

    sections = [
        {
            "title": "Platform & AI Layer", "color": ACCENT_PURPLE,
            "items": [
                ("NVIDIA", "Jetson Orin/Thor", "Universal AI compute standard"),
                ("OpenAI", "GPT-4V / VLA", "Figure AI backbone partner"),
                ("Google DeepMind", "RT-X / Gemini Robotics", "Open embodied AI research"),
                ("Microsoft", "Azure Robotics Platform", "Enterprise cloud infrastructure"),
                ("Meta AI", "PyTorch + Habitat Sim", "Open-source simulation tools"),
                ("Tesla", "Dojo + FSD Chip", "Proprietary, Optimus-only"),
            ],
            "x": 0.2
        },
        {
            "title": "OEM Leaders", "color": USA_BLUE,
            "items": [
                ("Tesla Optimus", "$20K target BOM", "22 DOF  |  <50 deployed  |  FSD chipset"),
                ("Boston Dynamics", "Atlas (electric)", "29 DOF  |  Hyundai owned  |  R&D focus"),
                ("Figure AI", "$675M raised", "Figure 02  |  BMW pilot  |  OpenAI VLA"),
                ("Agility Robotics", "Digit v5", "NASDAQ listed  |  Amazon warehouse"),
                ("1X Technologies", "NEO B2", "OpenAI backed  |  Norway/US"),
                ("Apptronik", "Apollo", "Mercedes-Benz manufacturing pilot"),
            ],
            "x": 3.45
        },
        {
            "title": "Component Suppliers", "color": ACCENT_CYAN,
            "items": [
                ("ATI Industrial Automation", "6-Axis F/T Sensors", "60%+ global market share"),
                ("Moog Inc.", "Servo Actuators / Drives", "Aerospace & defense grade"),
                ("Harmonic Drive LLC", "Precision Gearboxes", "High-end robot standard"),
                ("Intel RealSense", "D435/D455 Depth Cam", "Most US robots use this"),
                ("Qualcomm", "RB5/RB6 Platform", "Mid-range robot compute"),
                ("Synapticon", "Servo Drives", "Motion control specialist"),
            ],
            "x": 6.7
        }
    ]

    for sec in sections:
        sx = sec["x"]
        color = sec["color"]
        hb = add_rect(slide, sx, 1.45, 3.1, 0.28,
                      fill_rgb=RGBColor(int(color[0]*0.12), int(color[1]*0.12), int(color[2]*0.12)),
                      line_rgb=color, line_pt=0.75)
        add_text(slide, sec["title"], sx+0.1, 1.49, 2.9, 0.22, size=10.5, color=color, bold=True)

        for ii, (name, spec, note) in enumerate(sec["items"]):
            iy = 1.78 + ii * 0.58
            ib = add_rect(slide, sx, iy, 3.1, 0.54, fill_rgb=BG_CARD, line_rgb=RGBColor(0x30, 0x30, 0x50), line_pt=0.3)
            add_text(slide, name, sx+0.1, iy+0.04, 2.9, 0.22, size=10, color=TEXT_WHITE, bold=True)
            add_text(slide, spec, sx+0.1, iy+0.27, 1.8, 0.2, size=9, color=color)
            add_text(slide, note, sx+0.1, iy+0.38, 2.9, 0.18, size=8, color=TEXT_GRAY)

    # Bottom vulnerability banner
    vb = add_rect(slide, 0.2, 5.04, 9.6, 0.28, fill_rgb=RGBColor(0x28, 0x0A, 0x14), line_rgb=ACCENT_PINK, line_pt=0.75)
    add_text(slide,
             "Critical Vulnerability:  70% import dependency on key components — ATI force sensors, Harmonic Drive gearboxes, TSMC chips, CATL batteries. All US OEMs combined <100 units deployed in 2024.",
             0.35, 5.07, 9.3, 0.22, size=9, color=ACCENT_PINK)

    return slide


# ===================== SLIDE M5: China vs US Comparison =====================
def build_m5(prs):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    add_header(slide, "M - Mapping  China vs US Supply Chain",
               "Supply Chain Competitive Analysis",
               "7 Dimensions  |  Manufacturing Scale vs Platform Control  |  2027 Convergence",
               11, 27)

    dims = [
        ("Manufacturing Scale", "CN: Pearl River Delta cluster, 30+ suppliers, <$5K assembly\nUS: No domestic ecosystem, relies on Asian imports for mechanical parts", 95, 20),
        ("Actuator Ecosystem", "CN: Full stack harmonic/servo/F-T sensors, 30+ domestic\nUS: Moog/Harmonic Drive (high-end only), no mass production", 90, 30),
        ("AI & Compute Platform", "CN: Horizon/Black Sesame, 2-3yr behind CUDA ecosystem\nUS: NVIDIA monopoly: CUDA + Isaac Sim + Omniverse, 3-5yr moat", 35, 95),
        ("Deployment Volume", "CN: Unitree 5,500+ units; total industry >10,000 deployed\nUS: All OEMs combined <100 units; BMW pilot = 1 Figure robot", 98, 15),
        ("Cost Competitiveness", "CN: 65% of US baseline today; 35% at 100K (55% advantage)\nUS: $50K+ BOM; no path to $10K without China supply chain", 92, 100),
        ("Enterprise Adoption", "CN: Domestic auto/electronics pilots (BYD, CATL), China-only\nUS: BMW, Amazon, Mercedes global enterprise credibility", 50, 72),
        ("Geopolitical Risk", "CN: Low; fully domestic supply chain, minimal import dependency\nUS: High; ATI/TSMC/Harmonic concentration risk", 90, 28),
    ]

    add_text(slide, "CHINA", 0.25, 1.38, 3.35, 0.22, size=13, color=CHINA_RED, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, "DIMENSION", 3.65, 1.38, 2.7, 0.22, size=11, color=ACCENT_AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "US", 6.4, 1.38, 3.35, 0.22, size=13, color=USA_BLUE, bold=True, align=PP_ALIGN.LEFT)

    BAR_W = 3.15

    for i, (dim_name, notes, cn_score, us_score) in enumerate(dims):
        ry = 1.63 + i * 0.52
        row_bg = RGBColor(0x14, 0x14, 0x22) if i % 2 == 0 else RGBColor(0x10, 0x10, 0x1C)
        add_rect(slide, 0.2, ry, 9.6, 0.49, fill_rgb=row_bg)

        # Center label
        add_text(slide, dim_name, 3.65, ry+0.07, 2.7, 0.22, size=9.5, color=ACCENT_AMBER, bold=True, align=PP_ALIGN.CENTER)

        # China bar (right-to-left from center)
        cn_bar_w = BAR_W * (cn_score / 100.0)
        add_rect(slide, 0.25, ry+0.13, BAR_W, 0.22, fill_rgb=RGBColor(0x22, 0x22, 0x33))
        if cn_bar_w > 0:
            fill_x = 0.25 + (BAR_W - cn_bar_w)
            add_rect(slide, fill_x, ry+0.13, cn_bar_w, 0.22, fill_rgb=RGBColor(0xCC, 0x44, 0x44))
        add_text(slide, f"{cn_score}", 0.25, ry+0.15, BAR_W, 0.2, size=9, color=TEXT_WHITE, bold=True, align=PP_ALIGN.RIGHT)

        # US bar (left-to-right)
        us_bar_w = BAR_W * (us_score / 100.0)
        add_rect(slide, 6.55, ry+0.13, BAR_W, 0.22, fill_rgb=RGBColor(0x22, 0x22, 0x33))
        if us_bar_w > 0:
            add_rect(slide, 6.55, ry+0.13, us_bar_w, 0.22, fill_rgb=RGBColor(0x33, 0x66, 0xCC))
        add_text(slide, f"{us_score}", 6.55, ry+0.15, BAR_W, 0.2, size=9, color=TEXT_WHITE, bold=True, align=PP_ALIGN.LEFT)

    # Bottom insight
    bi = add_rect(slide, 0.2, 5.27, 9.6, 0.3, fill_rgb=BG_CARD, line_rgb=ACCENT_CYAN, line_pt=0.5)
    add_text(slide,
             "Key Insight: China dominates hardware manufacturing (9x deployment gap) while US controls AI platform layer. The 2026-2027 window is decisive — first mover to 100K units wins.",
             0.35, 5.31, 9.3, 0.24, size=9.5, color=TEXT_LGRAY)

    return slide


# ===================== SLIDE M6: Cost Curve =====================
def build_m6(prs):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    add_header(slide, "M - Mapping  Cost Curve Analysis",
               "Scale Economics & Cost Curve",
               "100K Units: China 35 vs US 100  |  65% Cost Advantage  |  Convergence to $10K Target",
               12, 27)

    # Scale comparison table (visual bar chart)
    scales = [
        ("Current", 65, 100, 95, 90),
        ("1K Units", 58, 100, 92, 88),
        ("10K Units", 45, 100, 88, 82),
        ("100K Units", 35, 100, 85, 78),
    ]

    cell_x = [0.3, 1.9, 3.7, 5.5, 7.3]
    headers = ["Scale", "China", "US (Base=100)", "Japan/EU", "Advantage"]
    h_colors = [TEXT_GRAY, CHINA_RED, USA_BLUE, EU_PURPLE, ACCENT_GREEN]

    # Header row
    for hi, (hdr, hc) in enumerate(zip(headers, h_colors)):
        add_text(slide, hdr, cell_x[hi], 1.45, 1.6, 0.28, size=10, color=hc, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.25, 1.7, 9.5, 0.015, fill_rgb=RGBColor(0x30, 0x30, 0x50))

    for ri, (scale, cn, us, jp, eu2) in enumerate(scales):
        ry = 1.75 + ri * 0.68
        row_bg = RGBColor(0x14, 0x14, 0x22) if ri % 2 == 0 else RGBColor(0x10, 0x10, 0x1C)
        add_rect(slide, 0.25, ry, 9.5, 0.65, fill_rgb=row_bg)

        # Scale label
        add_text(slide, scale, cell_x[0], ry+0.08, 1.5, 0.25, size=11, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Bar for each country
        bar_data = [(cell_x[1], cn, CHINA_RED), (cell_x[2], us, USA_BLUE), (cell_x[3], jp, EU_PURPLE)]
        for bx, val, bc in bar_data:
            bar_w = 1.5 * (val / 100.0)
            add_rect(slide, bx, ry+0.08, 1.5, 0.22, fill_rgb=RGBColor(0x22, 0x22, 0x33))
            if bar_w > 0:
                add_rect(slide, bx, ry+0.08, bar_w, 0.22, fill_rgb=bc)
            add_text(slide, f"{val}", bx, ry+0.08, 1.5, 0.22, size=10, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Advantage
        adv = 100 - cn
        adv_color = ACCENT_GREEN if adv > 50 else (ACCENT_AMBER if adv > 30 else TEXT_GRAY)
        add_text(slide, f"+{adv}%", cell_x[4], ry+0.08, 1.6, 0.22, size=14, color=adv_color, bold=True, align=PP_ALIGN.CENTER)

        # Notes row
        notes = [
            "Current prototypes\nSmall batch production",
            "Early adopter\nR&D deployment",
            "Industrial pilot\nEnterprise deployment",
            "Mass manufacturing\nConsumer market",
        ]
        add_text(slide, notes[ri], cell_x[0], ry+0.35, 9.5, 0.26, size=7.5, color=TEXT_GRAY, align=PP_ALIGN.LEFT)

    # Bottom: Analysis boxes
    panels = [
        ("55%", CHINA_RED, "Cost Advantage at 100K Scale", "35 (China) vs 100 (US). Structural manufacturing moat from vertical integration + labor + volume."),
        ("2027", ACCENT_AMBER, "Convergence Year", "China closes AI gap via data scale. US reshoring begins. Supply chain competition peaks."),
        ("$10K", ACCENT_CYAN, "Target Robot BOM", "China's path to $10K robot (from $25K today). Battery + actuators = 60% of BOM, heavily China-advantaged."),
    ]

    for pi, (stat, sc, title, desc) in enumerate(panels):
        px = 0.25 + pi * 3.2
        pb = add_rect(slide, px, 4.63, 3.05, 0.6, fill_rgb=BG_CARD, line_rgb=sc, line_pt=0.6)
        add_text(slide, stat, px+0.1, 4.66, 0.9, 0.32, size=20, color=sc, bold=True)
        add_text(slide, title, px+1.0, 4.67, 2.0, 0.2, size=9.5, color=TEXT_WHITE, bold=True)
        add_text(slide, desc, px+0.1, 4.9, 2.9, 0.32, size=8, color=TEXT_GRAY)

    return slide


# ===================== SLIDE M7: Convergence Timeline =====================
def build_m7(prs):
    slide = prs.slides.add_slide(blank)
    set_bg(slide, BG_DARK)
    add_header(slide, "M - Mapping  Supply Chain Timeline",
               "Convergence Timeline: 2023-2028+",
               "Technology Validation  Pilot Deployment  Critical Inflection  Winner Takes Most",
               13, 27)

    events = [
        ("2023-2024", "Technology Validation Phase  技术验证期", ACCENT_GREEN, "past", [
            "Both sides demonstrate working prototypes; supply chains remain largely separate",
            "China: Unitree G1 launched; UBTECH Walker S enters production; 3-5 home OEMs funded",
            "US: Boston Dynamics Atlas goes electric; Figure raises $75M; Agility IPOs on NASDAQ",
            "NVIDIA Jetson Orin becomes de-facto compute standard for US robotics ecosystem",
        ]),
        ("2025 (Now)", "Pilot Deployment Phase  试点部署期", ACCENT_AMBER, "current", [
            "Unitree ships 5,500+ units globally — ALL US OEMs combined still under 100 units",
            "China's manufacturing cost advantage becomes mathematically undeniable at scale",
            "US enterprise pilots: BMW (Figure 1 robot), Amazon (Agility Digit), DHL logistics",
            "Unitree IPO prospectus reveals truth: 73.6% customers = research/education, not industry",
        ]),
        ("2026-2027", "Critical Inflection  关键拐点", ACCENT_PINK, "future", [
            "China's AI platform gap narrows — domestic LLMs trained on 100M+ hours robot data",
            "Unitree targets 75,000 units/year capacity; AgiBot targets 10,000+ units",
            "US OEMs face manufacturing ceiling: Figure, Agility struggle to exceed 500 units/yr",
            "Geopolitical escalation: export controls vs domestic Chinese component substitution",
        ]),
        ("2028+", "Winner Takes Most  赢者通吃", ACCENT_PURPLE, "future", [
            "First company to 100K units creates insurmountable cost moat (35 vs 100 baseline)",
            "Supply chain regionalization solidifies; separate China / US-aligned ecosystems emerge",
            "Platform layer remains US-dominated (NVIDIA CUDA + simulation + AI training infra)",
            "Robot BOM reaches $10K-15K range; mass commercial and consumer adoption begins",
        ]),
    ]

    # Vertical line
    lx = 0.55
    add_rect(slide, lx+0.1, 1.42, 0.018, 3.72, fill_rgb=RGBColor(0x30, 0x30, 0x55))

    for ti, (date, title, color, phase, bullets) in enumerate(events):
        ty = 1.42 + ti * 0.97

        # Timeline dot
        dot = slide.shapes.add_shape(9, Inches(lx), Inches(ty+0.06), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        if phase == "current":
            dot.line.color.rgb = TEXT_WHITE
            dot.line.width = Pt(1.5)
        else:
            dot.line.fill.background()

        # Date
        add_text(slide, date, lx+0.35, ty, 1.35, 0.22, size=10, color=color, bold=True)
        # Title
        add_text(slide, title, lx+1.7, ty, 7.8, 0.22, size=10.5, color=TEXT_WHITE, bold=True)

        # Bullets
        for bi, bullet in enumerate(bullets):
            add_text(slide, f"  {chr(8226)}  {bullet}", lx+1.7, ty+0.25+bi*0.17, 7.8, 0.18, size=8.5, color=TEXT_LGRAY)

        # Separator
        if ti < len(events)-1:
            add_rect(slide, lx+0.35, ty+0.9, 9.2, 0.008, fill_rgb=RGBColor(0x28, 0x28, 0x40))

    # Bottom callout
    bc = add_rect(slide, 0.2, 5.2, 9.6, 0.3, fill_rgb=BG_CARD, line_rgb=ACCENT_CYAN, line_pt=0.5)
    add_text(slide,
             "Strategic Insight:  The 2026-2027 window is decisive. China leads hardware deployment 9:1; US controls AI intelligence layer. First mover to 100K units builds insurmountable supply chain moat.",
             0.35, 5.24, 9.3, 0.23, size=9.5, color=TEXT_LGRAY)

    return slide


# ===================== Main: Build all slides then insert at position 9 =====================
import copy
from pptx.oxml.ns import qn

def insert_slide_after(prs, new_slide, after_index):
    """Insert a slide into prs at position after_index+1 (0-indexed)"""
    xml_slides = prs.slides._sldIdLst
    # Get the new slide's rId
    new_rId = None
    for rel in prs.slides._sldIdLst.getparent().getparent().part.rels.values():
        if hasattr(rel, '_target') and hasattr(rel._target, 'slide') and False:
            pass
    # Simple approach: the new slide is currently last; move it to after_index+1
    slide_count = len(prs.slides)
    # Move from last position to target position
    target_pos = after_index + 1
    # The new slide is at position slide_count - 1; we need to move it to target_pos
    items = list(xml_slides)
    if len(items) > 0:
        last_item = items[-1]
        xml_slides.remove(last_item)
        # Insert at target position
        if target_pos >= len(list(xml_slides)):
            xml_slides.append(last_item)
        else:
            ref_item = list(xml_slides)[target_pos]
            xml_slides.insert(list(xml_slides).index(ref_item), last_item)


# Build 5 new slides (they get appended at the end first)
print("Building slide M3 (China Supply Chain)...")
m3 = build_m3(prs)
print("Building slide M4 (US Supply Chain)...")
m4 = build_m4(prs)
print("Building slide M5 (China vs US Comparison)...")
m5 = build_m5(prs)
print("Building slide M6 (Cost Curve)...")
m6 = build_m6(prs)
print("Building slide M7 (Timeline)...")
m7 = build_m7(prs)

# Now we have 27 slides (22 original + 5 new at the end)
# We need to move slides [22,23,24,25,26] (0-indexed) to after slide index 7 (slide 8 in 1-indexed)
# Strategy: move them one by one from end to target position
print(f"After adding: {len(prs.slides)} slides")

xml_slides = prs.slides._sldIdLst
all_items = list(xml_slides)
print(f"sldIdLst count: {len(all_items)}")

# We want order: [0..7] [22,23,24,25,26] [8..21]
# Current order: [0..21] [22,23,24,25,26]
# Extract the 5 new items (they're at positions 22-26)
new_items = all_items[22:27]
# Remove them from current positions
for item in new_items:
    xml_slides.remove(item)
# Insert them at position 8 (after index 7)
ref_item = list(xml_slides)[8]  # This is original slide 9 (0-indexed)
for item in reversed(new_items):
    ref_item = list(xml_slides)[8]
    xml_slides.insert(list(xml_slides).index(ref_item), item)

# Verify
print(f"Final slide count: {len(prs.slides)}")
all_items_final = list(xml_slides)
print(f"Final sldIdLst count: {len(all_items_final)}")

prs.save(OUTPUT_PATH)
print(f"Saved to: {OUTPUT_PATH}")

# Quick verification
prs2 = Presentation(OUTPUT_PATH)
print(f"Verification - total slides: {len(prs2.slides)}")
for i, slide in enumerate(prs2.slides):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            t = shape.text.strip()[:60].encode('ascii','replace').decode('ascii')
            texts.append(t)
    if texts:
        print(f"Slide {i+1}: {' | '.join(texts[:2])}")

