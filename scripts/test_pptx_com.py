import os
import comtypes.client

# Get PowerPoint version
ppt = comtypes.client.CreateObject("PowerPoint.Application")
print(f"PowerPoint version: {ppt.Version}")
print(f"Build: {ppt.Build}")
ppt.Quit()

# Create a minimal test PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.62)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Test Slide"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)

# Add content
body_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
tf = body_box.text_frame
p = tf.paragraphs[0]
p.text = "This is a test PPTX created by python-pptx"
p.font.size = Pt(18)

test_pptx = r"C:\Users\51229\WorkBuddy\Claw\test_minimal.pptx"
prs.save(test_pptx)
print(f"Test PPTX created: {os.path.getsize(test_pptx):,} bytes")

# Now try to open with PowerPoint COM
ppt2 = comtypes.client.CreateObject("PowerPoint.Application")
ppt2.Visible = 1
print("Opening test PPTX with PowerPoint COM...")
try:
    pres = ppt2.Presentations.Open(test_pptx, WithWindow=False)
    print("SUCCESS!")
    test_pdf = r"C:\Users\51229\WorkBuddy\Claw\test_minimal.pdf"
    pres.SaveAs(test_pdf, 32)
    print(f"PDF created: {os.path.exists(test_pdf)}")
    pres.Close()
except Exception as e:
    print(f"ERROR: {e}")
ppt2.Quit()
